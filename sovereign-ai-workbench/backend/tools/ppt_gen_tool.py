import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from backend.config import config

class PPTGenTool:
    """
    Real Presentation (.pptx) Generator Tool.
    Generates structured 6-slide management presentations.
    """

    def generate_presentation(self, title: str, slides_data: list, filename: str = "Management_Presentation.pptx") -> str:
        prs = Presentation()
        # 16:9 widescreen layout
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_slide_layout = prs.slide_layouts[6]

        # Slide 1: Title Slide
        slide1 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        p.font.name = "Arial"

        p2 = tf.add_paragraph()
        p2.text = "Sovereign Industrial AI Intelligence Platform | Air-Gapped Confidential Report"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(100, 116, 139)
        p2.font.name = "Arial"

        # Content Slides
        for slide_info in slides_data:
            s = prs.slides.add_slide(blank_slide_layout)
            
            # Slide Title Header
            header_box = s.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.0))
            tf_h = header_box.text_frame
            p_h = tf_h.paragraphs[0]
            p_h.text = slide_info.get("header", "Slide Header")
            p_h.font.size = Pt(24)
            p_h.font.bold = True
            p_h.font.color.rgb = RGBColor(30, 41, 59)
            p_h.font.name = "Arial"

            # Slide Content Body
            body_box = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0))
            tf_b = body_box.text_frame
            tf_b.word_wrap = True

            bullets = slide_info.get("bullets", [])
            for idx, b in enumerate(bullets):
                if idx == 0:
                    p_b = tf_b.paragraphs[0]
                else:
                    p_b = tf_b.add_paragraph()
                p_b.text = f"• {b}"
                p_b.font.size = Pt(16)
                p_b.font.color.rgb = RGBColor(51, 65, 85)
                p_b.font.name = "Arial"
                p_b.space_after = Pt(12)

        output_path = os.path.join(config.DELIVERABLES_DIR, filename)
        prs.save(output_path)
        return output_path

ppt_gen_tool = PPTGenTool()
